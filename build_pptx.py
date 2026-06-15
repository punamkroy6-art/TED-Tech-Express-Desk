"""
TED Business Case — python-pptx generator
Run: python build_pptx.py
Output: TED_Business_Case.pptx  (same directory as this script)
"""

import os, sys

# ── install if needed ────────────────────────────────────────────────────────
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy
from lxml import etree

# ── colour palette ────────────────────────────────────────────────────────────
NAVY   = "0D1F3C"
NAVY2  = "163360"
CYAN   = "0EA5E9"
CYAN2  = "BAE6FD"
WHITE  = "FFFFFF"
LIGHT  = "F8FAFC"
MUTED  = "64748B"
DARK   = "1E293B"
GREEN  = "16A34A"
AMBER  = "D97706"
BORDER = "E2E8F0"

# ── helpers ───────────────────────────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_bg(slide, hex_str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_str)

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_pt=0.75):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_hex is None:
        fill.background()
    else:
        fill.solid()
        fill.fore_color.rgb = rgb(fill_hex)
    ln = shape.line
    if line_hex:
        ln.color.rgb = rgb(line_hex)
        ln.width = Pt(line_pt)
    else:
        ln.fill.background()
    shape.shadow.inherit = False
    return shape

def add_oval(slide, x, y, w, h, fill_hex=None, line_hex=None, line_pt=1.5):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
    else:
        shape.fill.background()
    ln = shape.line
    if line_hex:
        ln.color.rgb = rgb(line_hex)
        ln.width = Pt(line_pt)
    else:
        ln.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, x, y, w, h, font_size,
             bold=False, color_hex=WHITE, align=PP_ALIGN.LEFT,
             italic=False, font_name="Calibri",
             wrap=True, margin_left=0.05, margin_top=0.02):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    # margins
    tf.margin_left  = Inches(margin_left)
    tf.margin_right = Inches(0.02)
    tf.margin_top   = Inches(margin_top)
    tf.margin_bottom= Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(font_size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = rgb(color_hex)
    f.name = font_name
    return txb

def add_text_multiline(slide, lines, x, y, w, h, font_size,
                        bold=False, color_hex=WHITE, align=PP_ALIGN.LEFT,
                        italic=False, font_name="Calibri"):
    """lines: list of (text, bold, color_hex, font_size, italic, font_name)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for (txt, b, c, fs, it, fn) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        f = run.font
        f.size = Pt(fs)
        f.bold = b
        f.italic = it
        f.color.rgb = rgb(c)
        f.name = fn
    return txb

# ── presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]  # completely blank

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, NAVY)

# top cyan bar
add_rect(s1, 0, 0, 10, 0.06, CYAN)

# eye icon — outer oval (no fill, CYAN border 2pt)
add_oval(s1, 1.15, 1.4, 1.1, 0.65, fill_hex=None, line_hex=CYAN, line_pt=2)
# pupil
add_oval(s1, 1.52, 1.57, 0.46, 0.46, fill_hex=CYAN, line_hex=None)
# inner dot
add_oval(s1, 1.65, 1.69, 0.20, 0.20, fill_hex=NAVY, line_hex=None)

# TED
add_text(s1, "TED", 1.1, 2.2, 2, 0.62, 38, bold=True, color_hex=WHITE,
         align=PP_ALIGN.LEFT, font_name="Calibri")
# TECH EXPRESS DESK
add_text(s1, "TECH EXPRESS DESK", 1.1, 2.82, 4.5, 0.32, 10,
         bold=False, color_hex=CYAN2, align=PP_ALIGN.LEFT, font_name="Calibri")
# Business Case
add_text(s1, "Business Case", 1.1, 3.22, 8, 0.68, 30, bold=True,
         color_hex=WHITE, align=PP_ALIGN.LEFT, font_name="Cambria")
# subtitle
add_text(s1, "AI-Powered Self-Service IT Support Kiosk  ·  Phase 1 Investment Proposal",
         1.1, 3.96, 8, 0.35, 12, bold=False, color_hex=MUTED,
         align=PP_ALIGN.LEFT, font_name="Calibri")
# cyan divider
add_rect(s1, 1.1, 4.44, 3.5, 0.05, CYAN)
# footer
add_text(s1, "IT Operations – AI/ML Team  ·  v1.0 Draft  ·  Confidential",
         1.1, 4.74, 7, 0.32, 9.5, bold=False, color_hex="4B5E78",
         align=PP_ALIGN.LEFT, font_name="Calibri")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, LIGHT)

add_text(s2, "THE PROBLEM", 0.5, 0.3, 3, 0.28, 8.5, bold=True,
         color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text(s2, "IT Support Is Broken for the Modern Workforce",
         0.5, 0.58, 9, 0.52, 24, bold=True, color_hex=DARK,
         align=PP_ALIGN.LEFT, font_name="Cambria")

pain_cards = [
    ("Slow Resolution",
     "Employees log a ticket and wait hours or days for simple, repeatable issues."),
    ("Repeat Tickets",
     "The same 20% of issue types account for 60-70% of all Service Desk volume."),
    ("No On-Site Help",
     "No walk-up intelligent support. Employees wait or self-diagnose unsuccessfully."),
    ("Thin Ticket Data",
     "Tickets arrive with minimal context, forcing SD agents to investigate from scratch."),
]
stat_boxes = [
    ("45 min",   "Avg. time to first SD response per ticket"),
    ("60-70%",   "Of tickets are repeat or low-complexity issues"),
    ("$22",      "Estimated cost per ticket resolved by SD agent"),
    ("2.3 days", "Avg. ticket resolution time across all priorities"),
]

base_y = 1.45
card_h = 0.82
gap    = 0.14

for i, (title, desc) in enumerate(pain_cards):
    cy = base_y + i * (card_h + gap)
    # white card
    add_rect(s2, 0.4, cy, 5.2, card_h, WHITE, line_hex=BORDER, line_pt=0.75)
    add_text(s2, title, 0.5, cy + 0.08, 5.0, 0.28, 10.5, bold=True,
             color_hex=DARK, align=PP_ALIGN.LEFT, font_name="Calibri")
    add_text(s2, desc,  0.5, cy + 0.38, 5.0, 0.38, 9, bold=False,
             color_hex=MUTED, align=PP_ALIGN.LEFT, font_name="Calibri")

    # stat box (NAVY)
    val, lbl = stat_boxes[i]
    sx = 5.9
    add_rect(s2, sx, cy, 3.65, card_h, NAVY)
    add_text(s2, val, sx + 0.08, cy + 0.08, 1.35, 0.50, 22, bold=True,
             color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Cambria")
    add_text(s2, lbl, sx + 1.5,  cy + 0.12, 2.0,  0.58, 8.5, bold=False,
             color_hex=CYAN2, align=PP_ALIGN.LEFT, font_name="Calibri")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW TED WORKS
# ═════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, WHITE)

add_text(s3, "THE SOLUTION", 0.5, 0.22, 3, 0.28, 8.5, bold=True,
         color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text(s3, "How TED Works — Three Intelligent Paths",
         0.5, 0.52, 9, 0.52, 24, bold=True, color_hex=DARK,
         align=PP_ALIGN.LEFT, font_name="Cambria")

# step circles
step_cx = [1.5, 4.55, 7.6]   # center x of circles
circle_r_x = 0.275            # half-width
circle_r_y = 0.275
steps = [
    ("Walk Up &\nAuthenticate", "Badge tap or SSO login"),
    ("Connect\nDevice",         "USB or manual issue select"),
    ("AI Diagnostics\nRun",     "Scan completes in under 60 sec"),
]
for i, cx in enumerate(step_cx):
    ox = cx - circle_r_x
    oy = 1.35
    add_oval(s3, ox, oy, 0.55, 0.55, fill_hex=CYAN, line_hex=None)
    add_text(s3, str(i+1), ox+0.12, oy+0.09, 0.31, 0.35, 14, bold=True,
             color_hex=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    # step title
    add_text(s3, steps[i][0], cx-0.8, 1.97, 1.6, 0.38, 10, bold=True,
             color_hex=DARK, align=PP_ALIGN.CENTER, font_name="Calibri")
    # step desc
    add_text(s3, steps[i][1], cx-0.8, 2.38, 1.6, 0.28, 8.5, bold=False,
             color_hex=MUTED, align=PP_ALIGN.CENTER, font_name="Calibri")

# arrows between circles
# circle right edge → next circle left edge
for i in range(len(step_cx)-1):
    ax = step_cx[i] + circle_r_x + 0.04
    bx = step_cx[i+1] - circle_r_x - 0.04
    aw = bx - ax
    add_rect(s3, ax, 1.35 + circle_r_y - 0.02, aw, 0.04, CYAN)

# divider
add_rect(s3, 0.4, 2.82, 9.2, 0.03, BORDER)

# path cards
path_cards = [
    ("Path A", GREEN, "Self-Service\nResolution",
     ["AI auto-fix applied", "Employee confirms fix worked", "Session ends — no ticket"],
     "~30% of sessions resolved"),
    ("Path B", CYAN, "Guided Fix +\nSD Ticket",
     ["AI provides step-by-step fix", "Screenshot & diagnosis captured", "Freshworks ticket auto-created"],
     "Pre-filled AI-enriched ticket"),
    ("Path C", AMBER, "Loaner Device\nDispensed",
     ["Hardware failure confirmed by AI", "Kiosk locker unlocks loaner", "Asset logged in Freshworks"],
     "Zero downtime for employee"),
]
card_xs = [0.4, 3.52, 6.64]
for i, (badge, color, title, bullets, outcome) in enumerate(path_cards):
    cx = card_xs[i]
    cy = 2.95
    cw = 2.95
    ch = 2.5
    add_rect(s3, cx, cy, cw, ch, LIGHT, line_hex=BORDER, line_pt=0.75)
    # badge
    add_rect(s3, cx+0.1, cy+0.1, 0.52, 0.22, color)
    add_text(s3, badge, cx+0.11, cy+0.10, 0.5, 0.22, 7.5, bold=True,
             color_hex=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    # title
    add_text(s3, title, cx+0.08, cy+0.38, cw-0.16, 0.46, 11, bold=True,
             color_hex=DARK, align=PP_ALIGN.LEFT, font_name="Calibri")
    # bullets
    for j, b in enumerate(bullets):
        add_text(s3, "•  " + b, cx+0.12, cy+0.9+j*0.34, cw-0.2, 0.30, 8.5,
                 bold=False, color_hex=MUTED, align=PP_ALIGN.LEFT, font_name="Calibri")
    # outcome
    add_rect(s3, cx+0.08, cy+2.12, cw-0.16, 0.28, color)
    add_text(s3, outcome, cx+0.10, cy+2.13, cw-0.2, 0.26, 8, bold=True,
             color_hex=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BUSINESS VALUE
# ═════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, NAVY)

add_text(s4, "BUSINESS VALUE", 0.5, 0.22, 4, 0.28, 8.5, bold=True,
         color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text(s4, "Measurable Impact from Day One",
         0.5, 0.52, 9, 0.52, 24, bold=True, color_hex=WHITE,
         align=PP_ALIGN.LEFT, font_name="Cambria")

kpis = [
    (">30%",   "Ticket Deflection Rate",    "Sessions resolved without any SD ticket"),
    (">25%",   "Self-Service Resolution",   "Issues resolved at kiosk — zero SD involvement"),
    ("<5 min", "Avg. Kiosk Session",        "From login to resolution or ticket creation"),
    (">80%",   "AI Diagnosis Accuracy",     "Issue category correct vs SD final outcome"),
    ("20%",    "Faster SD Resolution",      "Handle time vs standard tickets"),
    (">4.0",   "Employee CSAT Score",       "Post-session survey target out of 5.0"),
]
kpi_xs = [0.4, 3.48, 6.56]
kpi_ys = [1.42, 3.30]
for idx, (val, lbl, sub) in enumerate(kpis):
    row = idx // 3
    col = idx % 3
    kx = kpi_xs[col]
    ky = kpi_ys[row]
    kw = 2.88
    kh = 1.65
    add_rect(s4, kx, ky, kw, kh, NAVY2, line_hex="1E4DB7", line_pt=1)
    add_text(s4, val, kx+0.15, ky+0.15, kw-0.3, 0.55, 30, bold=True,
             color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Cambria")
    add_text(s4, lbl, kx+0.15, ky+0.72, kw-0.3, 0.32, 10, bold=True,
             color_hex=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
    add_text(s4, sub, kx+0.15, ky+1.06, kw-0.3, 0.48, 8.5, bold=False,
             color_hex=MUTED, align=PP_ALIGN.LEFT, font_name="Calibri")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROADMAP
# ═════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, LIGHT)

add_text(s5, "IMPLEMENTATION ROADMAP", 0.5, 0.22, 5, 0.28, 8.5, bold=True,
         color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text(s5, "5 Phases  ·  20 Weeks to Full Rollout",
         0.5, 0.52, 9, 0.52, 24, bold=True, color_hex=DARK,
         align=PP_ALIGN.LEFT, font_name="Cambria")

# timeline bar
add_rect(s5, 0.38, 1.5, 9.24, 0.07, BORDER)

phases = [
    ("Foundation",    "Wk 1-3",   "0369A1",
     ["Docker infra + CI/CD", "PostgreSQL + Redis", "SSO + Badge Auth", "API scaffolding"]),
    ("AI Core",       "Wk 4-7",   "0891B2",
     ["Error DB (100+ items)", "OCR pipeline", "LLM (Groq/Azure)", "RAG + pgvector"]),
    ("ITSM",          "Wk 8-10",  "0EA5E9",
     ["Freshworks integration", "Auto-ticket creation", "KB + ticket lookup", "Asset + loaner SR"]),
    ("Kiosk UI + QA", "Wk 11-14", "0EA5E9",
     ["React 18 touch UI", "All screen flows", "Security testing", "OWASP ZAP scan"]),
    ("Pilot + Rollout","Wk 15-20","38BDF8",
     ["1-2 kiosk pilot", "Grafana monitoring", "CSAT tracking", "Full site rollout"]),
]
card_xs5 = [0.42, 2.29, 4.16, 6.03, 7.90]
for i, (name, weeks, color, items) in enumerate(phases):
    px = card_xs5[i]
    py = 1.82
    pw = 1.76
    ph = 3.35
    # dot on timeline
    dot_x = px + 0.66
    add_oval(s5, dot_x, 1.37, 0.28, 0.28, fill_hex=color, line_hex=None)
    # card
    add_rect(s5, px, py, pw, ph, WHITE, line_hex=BORDER, line_pt=0.75)
    # phase number badge
    add_rect(s5, px+0.1, py+0.1, 0.38, 0.22, color)
    add_text(s5, f"0{i+1}", px+0.11, py+0.10, 0.36, 0.22, 8, bold=True,
             color_hex=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    # name
    add_text(s5, name, px+0.08, py+0.40, pw-0.16, 0.30, 10, bold=True,
             color_hex=DARK, align=PP_ALIGN.LEFT, font_name="Calibri")
    # weeks
    add_text(s5, weeks, px+0.08, py+0.72, pw-0.16, 0.24, 8.5, bold=True,
             color_hex=color, align=PP_ALIGN.LEFT, font_name="Calibri", italic=True)
    # deliverables
    for j, item in enumerate(items):
        add_text(s5, "• " + item, px+0.08, py+1.02+j*0.54, pw-0.12, 0.48, 8,
                 bold=False, color_hex=MUTED, align=PP_ALIGN.LEFT, font_name="Calibri")

# go-live footer
add_rect(s5, 0.38, 5.22, 9.24, 0.30, NAVY)
add_text(s5,
         "Go-Live Target: Week 20  ·  Full production rollout with live monitoring & deflection reporting",
         0.5, 5.24, 9.0, 0.26, 9, bold=False, color_hex=CYAN2,
         align=PP_ALIGN.CENTER, font_name="Calibri")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — NEXT STEPS
# ═════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, NAVY)

# bottom cyan bar
add_rect(s6, 0, 5.57, 10, 0.06, CYAN)

add_text(s6, "NEXT STEPS", 0.5, 0.22, 3, 0.28, 8.5, bold=True,
         color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text(s6, "What We Need to Proceed",
         0.5, 0.52, 9, 0.52, 26, bold=True, color_hex=WHITE,
         align=PP_ALIGN.LEFT, font_name="Cambria")

# column headers
add_text(s6, "Decisions Required", 0.5, 1.32, 4.4, 0.30, 11, bold=True,
         color_hex=CYAN2, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text(s6, "Immediate Actions", 5.3, 1.32, 4.2, 0.30, 11, bold=True,
         color_hex=CYAN2, align=PP_ALIGN.LEFT, font_name="Calibri")

decisions = [
    ("IT Manager",        "Phase 1 kiosk unit count and office site selection"),
    ("AI/ML Engineering", "LLM provider: Groq (cost) vs Azure OpenAI (SLA)"),
    ("IT Infrastructure", "SSO identity provider: Azure AD or Okta"),
    ("Asset Management",  "Loaner device inventory + approval workflow"),
    ("Legal/Compliance",  "Data retention policy for session recordings"),
    ("IT Operations",     "Freshworks API tier and rate limit confirmation"),
]
actions = [
    ("Week 1", "Approve Phase 1 budget & resource allocation"),
    ("Week 1", "Assign project lead and core engineering team"),
    ("Week 1", "Confirm Freshworks API access and credentials"),
    ("Week 2", "Complete site assessment for kiosk placement"),
    ("Week 2", "Kick off Phase 1: repo, infra, CI/CD"),
    ("Week 3", "Phase 1 checkpoint — proceed to AI Core"),
]

base_y6 = 1.74
card_h6 = 0.48
gap6    = 0.08

for i, (owner, item) in enumerate(decisions):
    cy = base_y6 + i * (card_h6 + gap6)
    add_rect(s6, 0.4, cy, 4.55, card_h6, NAVY2, line_hex="1E4DB7", line_pt=0.75)
    add_text(s6, owner, 0.52, cy+0.04, 4.3, 0.18, 7.5, bold=True,
             color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
    add_text(s6, item,  0.52, cy+0.23, 4.3, 0.22, 8.5, bold=False,
             color_hex="CADCFC", align=PP_ALIGN.LEFT, font_name="Calibri")

for i, (week, action) in enumerate(actions):
    cy = base_y6 + i * (card_h6 + gap6)
    add_rect(s6, 5.22, cy, 4.30, card_h6, "0C1F3A", line_hex=CYAN, line_pt=0.75)
    add_text(s6, week,   5.34, cy+0.04, 4.1, 0.18, 7.5, bold=True,
             color_hex=CYAN, align=PP_ALIGN.LEFT, font_name="Calibri")
    add_text(s6, action, 5.34, cy+0.23, 4.1, 0.22, 8.5, bold=False,
             color_hex=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")

# ── save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TED_Business_Case.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
