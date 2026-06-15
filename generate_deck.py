"""
TED Business Case — 6-Slide Deck Generator
Run: python generate_deck.py
Output: TED_Business_Case.pptx (saved next to this script)
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as I, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy, os

# ── helpers ──────────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width=0, radius=False):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1 if not radius else 5,  # MSO_SHAPE_TYPE: RECTANGLE=1, ROUNDED_RECT=5
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width) if line_width else Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_hex):
    shape = slide.shapes.add_shape(
        9,  # OVAL
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color_hex, width_pt=1.0):
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color_hex)
    connector.line.width = Pt(width_pt)
    return connector

def add_text(slide, text, x, y, w, h, font_size, bold=False, color_hex="1E293B",
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)
    run.font.name = font_name
    return txBox

def add_multiline(slide, lines, x, y, w, h, font_size, color_hex="1E293B",
                  font_name="Calibri", bold=False, line_spacing_pt=None):
    """lines = list of (text, color_hex, bold) tuples, or just strings"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, col, b = line, color_hex, bold
        else:
            txt, col, b = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = b
        run.font.color.rgb = rgb(col)
        run.font.name = font_name
    return txBox

# ── palette ──────────────────────────────────────────────────
NAVY   = "0D1F3C"
NAVY2  = "163360"
CYAN   = "0EA5E9"
CYAN2  = "BAE6FD"
WHITE  = "FFFFFF"
LIGHT  = "F8FAFC"
MUTED  = "64748B"
DARK   = "1E293B"
GREEN  = "16A34A"
AMBER  = "D97706"
BORDER = "E2E8F0"

# ── Presentation setup ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]  # completely blank


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY)

# Top accent line
add_rect(s, 0, 0, 10, 0.06, CYAN)

# Eye icon - outer oval
add_oval(s, 1.15, 1.4, 1.1, 0.65, NAVY2)
# Eye border effect via line
eye = s.shapes.add_shape(9, Inches(1.15), Inches(1.4), Inches(1.1), Inches(0.65))
eye.fill.background()
eye.line.color.rgb = rgb(CYAN)
eye.line.width = Pt(2)
# Pupil
add_oval(s, 1.52, 1.57, 0.46, 0.46, CYAN)
add_oval(s, 1.65, 1.69, 0.2, 0.2, NAVY)

# TED title
add_text(s, "TED", 1.1, 2.2, 2, 0.62, 38, bold=True, color_hex=WHITE, font_name="Calibri")
add_text(s, "TECH EXPRESS DESK", 1.1, 2.82, 4.5, 0.32, 10, color_hex=CYAN2, font_name="Calibri")

add_text(s, "Business Case", 1.1, 3.22, 8, 0.68, 30, bold=True, color_hex=WHITE, font_name="Cambria")
add_text(s, "AI-Powered Self-Service IT Support Kiosk  ·  Phase 1 Investment Proposal",
         1.1, 3.96, 8, 0.35, 12, color_hex=MUTED, font_name="Calibri")

# Divider
add_rect(s, 1.1, 4.44, 3.5, 0.05, CYAN)

add_text(s, "IT Operations – AI/ML Team  ·  v1.0 Draft  ·  Confidential",
         1.1, 4.74, 7, 0.32, 9.5, color_hex="4B5E78", font_name="Calibri")


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = rgb(LIGHT)

add_text(s, "THE PROBLEM", 0.5, 0.3, 4, 0.28, 8.5, bold=True, color_hex=CYAN, font_name="Calibri")
add_text(s, "IT Support Is Broken for the Modern Workforce",
         0.5, 0.58, 9, 0.62, 24, bold=True, color_hex=DARK, font_name="Cambria")

# Pain point cards (left)
pains = [
    ("Slow Resolution",
     "Employees log a ticket and wait hours or days for simple, repeatable issues."),
    ("Repeat Tickets",
     "The same 20% of issue types account for 60-70% of all Service Desk volume."),
    ("No On-Site Help",
     "No walk-up intelligent support. Employees wait or self-diagnose unsuccessfully."),
    ("Thin Ticket Data",
     "Tickets arrive with minimal context, forcing SD agents to investigate from scratch."),
]
for i, (title, desc) in enumerate(pains):
    y = 1.45 + i * 0.96
    add_rect(s, 0.4, y, 5.2, 0.82, WHITE, BORDER, 0.75)
    add_text(s, title, 0.55, y + 0.08, 4.8, 0.3, 10.5, bold=True, color_hex=DARK)
    add_text(s, desc,  0.55, y + 0.38, 4.8, 0.38, 9, color_hex=MUTED)

# Impact stats (right)
stats = [
    ("45 min",  "Avg. time to first\nSD response per ticket"),
    ("60-70%",  "Of tickets are repeat\nor low-complexity issues"),
    ("$22",     "Estimated cost per ticket\nresolved by SD agent"),
    ("2.3 days","Avg. ticket resolution\ntime across all priorities"),
]
for i, (val, label) in enumerate(stats):
    y = 1.45 + i * 0.96
    add_rect(s, 5.9, y, 3.65, 0.82, NAVY)
    add_text(s, val,   5.98, y + 0.08, 1.5, 0.54, 22, bold=True, color_hex=CYAN, font_name="Cambria")
    add_text(s, label, 7.45, y + 0.12, 2.0, 0.58, 8.5, color_hex=CYAN2)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW TED WORKS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = rgb(WHITE)

add_text(s, "THE SOLUTION", 0.5, 0.3, 4, 0.28, 8.5, bold=True, color_hex=CYAN)
add_text(s, "How TED Works — Three Intelligent Paths",
         0.5, 0.58, 9, 0.58, 24, bold=True, color_hex=DARK, font_name="Cambria")

# Entry steps (3 circles + arrows)
steps = [
    ("1", "Walk Up & Authenticate", "Badge tap or SSO login"),
    ("2", "Connect Device",         "USB or manual issue select"),
    ("3", "AI Diagnostics Run",     "Scan completes in <60 sec"),
]
for i, (num, title, desc) in enumerate(steps):
    cx = 0.65 + i * 3.05
    add_oval(s, cx + 0.85, 1.32, 0.55, 0.55, CYAN)
    add_text(s, num, cx + 0.85, 1.32, 0.55, 0.55, 14, bold=True, color_hex=WHITE,
             align=PP_ALIGN.CENTER)
    if i < 2:
        add_rect(s, cx + 2.6, 1.53, 0.55, 0.04, CYAN)
    add_text(s, title, cx + 0.32, 1.95, 2.5, 0.32, 10, bold=True, color_hex=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, desc,  cx + 0.32, 2.28, 2.5, 0.32, 8.5, color_hex=MUTED,
             align=PP_ALIGN.CENTER)

# Divider
add_rect(s, 0.4, 2.78, 9.2, 0.04, BORDER)

# Three path cards
paths = [
    (GREEN, "Path A", "Self-Service Resolution",
     ["AI auto-fix applied", "Employee confirms fix worked", "Session ends — no ticket created"],
     "~30% of sessions resolved"),
    (CYAN,  "Path B", "Guided Fix + SD Ticket",
     ["AI provides step-by-step fix", "Screenshot & diagnosis captured", "Freshworks ticket auto-created"],
     "Pre-filled, AI-enriched ticket"),
    (AMBER, "Path C", "Loaner Device Dispensed",
     ["Hardware failure confirmed by AI", "Kiosk locker unlocks loaner", "Asset logged in Freshworks"],
     "Zero downtime for employee"),
]
for i, (color, label, head, items, outcome) in enumerate(paths):
    x = 0.4 + i * 3.12
    y = 2.9
    add_rect(s, x, y, 2.95, 2.52, LIGHT, BORDER, 0.75)
    # Label badge
    add_rect(s, x + 0.12, y + 0.12, 0.65, 0.26, color)
    add_text(s, label, x + 0.12, y + 0.12, 0.65, 0.26, 8, bold=True, color_hex=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, head, x + 0.12, y + 0.46, 2.7, 0.35, 10.5, bold=True, color_hex=DARK)
    for li, item in enumerate(items):
        add_text(s, f"→  {item}", x + 0.12, y + 0.87 + li * 0.38, 2.7, 0.34,
                 8.5, color_hex=MUTED)
    add_text(s, outcome, x + 0.12, y + 2.18, 2.7, 0.25, 8, bold=True,
             color_hex=color, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — BUSINESS VALUE (dark)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY)

add_text(s, "BUSINESS VALUE", 0.5, 0.3, 5, 0.28, 8.5, bold=True, color_hex=CYAN)
add_text(s, "Measurable Impact from Day One",
         0.5, 0.58, 9, 0.6, 24, bold=True, color_hex=WHITE, font_name="Cambria")

kpis = [
    (">30%",   "Ticket Deflection Rate",      "Sessions resolved without any SD ticket"),
    (">25%",   "Self-Service Resolution",     "Issues resolved at kiosk — zero SD involvement"),
    ("<5 min", "Avg. Kiosk Session",          "From login to resolution or ticket creation"),
    (">80%",   "AI Diagnosis Accuracy",       "Issue category correct vs SD final outcome"),
    ("20%",    "Faster SD Resolution",        "Handle time vs standard tickets"),
    (">4.0",   "Employee CSAT Score",         "Post-session survey target out of 5.0"),
]
for i, (val, label, sub) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 3.08
    y = 1.42 + row * 1.88
    add_rect(s, x, y, 2.88, 1.65, NAVY2, "1E4DB7", 0.75)
    add_text(s, val,   x + 0.18, y + 0.16, 2.55, 0.58, 30, bold=True,
             color_hex=CYAN, font_name="Cambria")
    add_text(s, label, x + 0.18, y + 0.76, 2.55, 0.32, 10, bold=True, color_hex=WHITE)
    add_text(s, sub,   x + 0.18, y + 1.08, 2.55, 0.46, 8.5, color_hex=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = rgb(LIGHT)

add_text(s, "IMPLEMENTATION ROADMAP", 0.5, 0.3, 6, 0.28, 8.5, bold=True, color_hex=CYAN)
add_text(s, "5 Phases · 20 Weeks to Full Rollout",
         0.5, 0.58, 9, 0.58, 24, bold=True, color_hex=DARK, font_name="Cambria")

# Timeline bar
add_rect(s, 0.38, 1.5, 9.24, 0.07, BORDER)

phases = [
    ("01", "Foundation",     "Wk 1-3",  "0369A1",
     ["Docker infra + CI/CD", "PostgreSQL + Redis", "SSO + Badge Auth", "API scaffolding"]),
    ("02", "AI Core",        "Wk 4-7",  "0891B2",
     ["Error DB (100+ items)", "OCR pipeline", "LLM (Groq/Azure)", "RAG + pgvector"]),
    ("03", "ITSM",           "Wk 8-10", "0EA5E9",
     ["Freshworks integration", "Auto-ticket creation", "KB + ticket lookup", "Asset + loaner SR"]),
    ("04", "Kiosk UI + QA",  "Wk 11-14","0EA5E9",
     ["React 18 touch UI", "All screen flows", "Security testing", "OWASP ZAP scan"]),
    ("05", "Pilot + Rollout","Wk 15-20","38BDF8",
     ["1-2 kiosk pilot", "Grafana monitoring", "CSAT tracking", "Full site rollout"]),
]

for i, (num, name, weeks, color, items) in enumerate(phases):
    x = 0.38 + i * 1.87
    y_card = 1.82

    # Dot on timeline
    add_oval(s, x + 0.62, 1.37, 0.28, 0.28, color)

    # Card
    add_rect(s, x + 0.04, y_card, 1.76, 3.35, WHITE, BORDER, 0.75)

    # Phase badge
    add_rect(s, x + 0.14, y_card + 0.1, 0.44, 0.28, color)
    add_text(s, num, x + 0.14, y_card + 0.1, 0.44, 0.28, 9, bold=True,
             color_hex=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, name, x + 0.1, y_card + 0.46, 1.6, 0.38, 10, bold=True, color_hex=DARK)
    add_text(s, weeks, x + 0.1, y_card + 0.84, 1.6, 0.26, 8.5, bold=True,
             color_hex=color, italic=True)

    for li, item in enumerate(items):
        add_text(s, f"·  {item}", x + 0.1, y_card + 1.18 + li * 0.5, 1.62, 0.44,
                 8, color_hex=MUTED)

# Go-live footer
add_rect(s, 0.38, 5.22, 9.24, 0.3, NAVY)
add_text(s, "Go-Live Target: Week 20  ·  Full production rollout with live monitoring & deflection reporting",
         0.5, 5.22, 9.0, 0.3, 9, color_hex=CYAN2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — NEXT STEPS (dark)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY)
add_rect(s, 0, 5.57, 10, 0.06, CYAN)

add_text(s, "NEXT STEPS", 0.5, 0.3, 4, 0.28, 8.5, bold=True, color_hex=CYAN)
add_text(s, "What We Need to Proceed",
         0.5, 0.58, 9, 0.6, 26, bold=True, color_hex=WHITE, font_name="Cambria")

add_text(s, "Decisions Required", 0.5, 1.32, 4.5, 0.32, 11, bold=True, color_hex=CYAN2)

decisions = [
    ("IT Manager",       "Phase 1 kiosk unit count and office site selection"),
    ("AI/ML Engineering","LLM provider: Groq (cost) vs Azure OpenAI (SLA)"),
    ("IT Infrastructure","SSO identity provider: Azure AD or Okta"),
    ("Asset Management", "Loaner device inventory + approval workflow"),
    ("Legal/Compliance", "Data retention policy for session recordings"),
    ("IT Operations",    "Freshworks API tier and rate limit confirmation"),
]
for i, (owner, item) in enumerate(decisions):
    y = 1.74 + i * 0.56
    add_rect(s, 0.4, y, 4.55, 0.48, NAVY2, "1E4DB7", 0.75)
    add_text(s, owner, 0.55, y + 0.04, 4.2, 0.2, 7.5, bold=True, color_hex=CYAN)
    add_text(s, item,  0.55, y + 0.23, 4.2, 0.2, 8.5, color_hex="CADCFC")

add_text(s, "Immediate Actions", 5.3, 1.32, 4.2, 0.32, 11, bold=True, color_hex=CYAN2)

actions = [
    ("Week 1", "Approve Phase 1 budget & resource allocation"),
    ("Week 1", "Assign project lead and core engineering team"),
    ("Week 1", "Confirm Freshworks API access and credentials"),
    ("Week 2", "Complete site assessment for kiosk placement"),
    ("Week 2", "Kick off Phase 1: repo, infra, CI/CD"),
    ("Week 3", "Phase 1 checkpoint — proceed to AI Core"),
]
for i, (week, action) in enumerate(actions):
    y = 1.74 + i * 0.56
    add_rect(s, 5.22, y, 4.3, 0.48, "0C1F3A", CYAN, 0.75)
    add_text(s, week,   5.35, y + 0.04, 0.9, 0.2, 7.5, bold=True, color_hex=CYAN)
    add_text(s, action, 5.35, y + 0.23, 4.0, 0.2, 8.5, color_hex=WHITE)


# ── Save ──────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TED_Business_Case.pptx")
prs.save(out)
print(f"✓  Saved: {out}")
